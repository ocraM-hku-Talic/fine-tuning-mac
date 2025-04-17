#!/usr/bin/env python3
"""
Course Material Q&A Data Generator for Model Fine-Tuning

This script extracts content from PowerPoint slides, Word documents, PDF files, and text files
and generates question-answer pairs for fine-tuning language models on course-specific content.
"""

import os
import json
import random
import argparse
import requests
from typing import List, Dict, Any
import jsonlines
from transformers import pipeline
import pptx  # python-pptx library for PowerPoint extraction
import docx  # python-docx library for Word document extraction
import re    # For text processing
try:
    import fitz  # PyMuPDF for PDF processing
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False
    print("Warning: PyMuPDF not installed. PDF support will be disabled.")
    print("To enable PDF support, install PyMuPDF: pip install pymupdf")

# Define question templates for educational content
QUESTION_TEMPLATES = [
    "What is {}?",
    "Can you explain {} in detail?",
    "How does {} work?",
    "What are the main components of {}?",
    "Why is {} important in this context?",
    "What's the difference between {} and {}?",
    "What are the advantages of {}?",
    "What are the applications of {}?",
    "How would you implement {}?",
    "Describe the process of {}.",
    "What challenges are associated with {}?",
    "In what scenarios would you use {}?",
    "Can you provide examples of {}?",
    "How does {} relate to {}?",
    "What are the key takeaways about {}?",
]

def parse_args():
    """Parse command line arguments."""
    parser = argparse.ArgumentParser(description="Generate Q&A pairs from course materials for LLM fine-tuning")
    
    parser.add_argument(
        "--slides_dir",
        type=str,
        required=True,
        help="Directory containing course materials (.pptx, .docx, .pdf, .txt files)"
    )
    parser.add_argument(
        "--output",
        type=str,
        default="../data/raw/course_qa_data.jsonl",
        help="Output file path (default: ../data/raw/course_qa_data.jsonl)"
    )
    parser.add_argument(
        "--num_samples",
        type=int,
        default=100,
        help="Target number of Q&A pairs to generate (default: 100)"
    )
    parser.add_argument(
        "--format",
        type=str,
        choices=["alpaca", "sharegpt", "finetune"],
        default="finetune",
        help="Dataset format (default: finetune for direct prompt/completion pairs)"
    )
    parser.add_argument(
        "--use_ollama",
        action="store_true",
        help="Use Ollama for generating Q&A pairs (recommended)"
    )
    parser.add_argument(
        "--ollama_model",
        type=str,
        default="deepseek-r1:latest ",
        help="Ollama model name to use for generation (default: deepseek-r1:latest )"
    )
    parser.add_argument(
        "--split",
        action="store_true",
        help="Split the generated data into train/validation/test sets"
    )
    parser.add_argument(
        "--batch_size",
        type=int,
        default=5,
        help="Number of content items to process per batch (default: 5)"
    )
    
    return parser.parse_args()

def extract_slide_content(ppt_file):
    """Extract text content from a PowerPoint file."""
    presentation = pptx.Presentation(ppt_file)
    slides_content = []
    
    for slide_num, slide in enumerate(presentation.slides):
        slide_text = ""
        
        # Extract slide title
        if slide.shapes.title and slide.shapes.title.text:
            slide_text += f"Title: {slide.shapes.title.text}\n\n"
        
        # Extract text from all shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                # Skip if it's the title we already added
                if shape == slide.shapes.title:
                    continue
                slide_text += f"{shape.text}\n"
        
        if slide_text.strip():
            slides_content.append({
                "slide_num": slide_num + 1,
                "content": slide_text.strip(),
                "file": os.path.basename(ppt_file),
                "type": "slide"
            })
    
    return slides_content

def extract_word_content(docx_file):
    """Extract text content from a Word document file."""
    doc = docx.Document(docx_file)
    content = []
    
    # Process each paragraph
    for para_num, para in enumerate(doc.paragraphs):
        if para.text.strip():  # Skip empty paragraphs
            content.append({
                "para_num": para_num + 1,
                "content": para.text.strip(),
                "file": os.path.basename(docx_file),
                "type": "word"
            })
    
    return content

def extract_text_content(txt_file):
    """Extract text content from a plain text file."""
    content = []
    
    with open(txt_file, 'r', encoding='utf-8', errors='replace') as f:
        lines = f.readlines()
        
        # Group lines into paragraphs (simple heuristic: empty lines separate paragraphs)
        paragraphs = []
        current_para = []
        
        for line in lines:
            if line.strip():
                current_para.append(line.strip())
            elif current_para:  # Empty line and we have content
                paragraphs.append(" ".join(current_para))
                current_para = []
        
        # Don't forget last paragraph if file doesn't end with empty line
        if current_para:
            paragraphs.append(" ".join(current_para))
        
        # Create content entries
        for para_num, para in enumerate(paragraphs):
            if para.strip():
                content.append({
                    "para_num": para_num + 1,
                    "content": para.strip(),
                    "file": os.path.basename(txt_file),
                    "type": "text"
                })
    
    return content

def extract_pdf_content(pdf_file):
    """Extract text content from a PDF file."""
    if not PDF_SUPPORT:
        print(f"Skipping PDF {pdf_file} as PyMuPDF is not installed")
        return []
    
    content = []
    try:
        doc = fitz.open(pdf_file)
        
        for page_num, page in enumerate(doc):
            # Extract text from the page
            text = page.get_text()
            if text.strip():
                # Split text into paragraphs based on newlines
                paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
                
                # Process each paragraph
                for para_num, para in enumerate(paragraphs):
                    if len(para) > 10:  # Skip very short paragraphs
                        content_item = {
                            "page_num": page_num + 1,
                            "para_num": para_num + 1,
                            "content": para.strip(),
                            "file": os.path.basename(pdf_file),
                            "type": "pdf"
                        }
                        content.append(content_item)
        
    except Exception as e:
        print(f"Error processing PDF {pdf_file}: {e}")
    
    return content

def extract_topics_from_content(content_items):
    """Extract potential topics from all content types."""
    topics = []
    
    for item in content_items:
        content = item["content"]
        
        # For PowerPoint slides, extract title if it exists
        if "slide_num" in item:
            title_match = content.split("\n\n")[0] if "\n\n" in content else ""
            if title_match.startswith("Title:"):
                topic = title_match[6:].strip()  # Remove "Title: " prefix
                if topic:
                    topics.append(topic)
        
        # Extract bullet points and short sentences (for all content types)
        lines = content.split("\n")
        for line in lines:
            line = line.strip()
            # Skip very short or very long lines, and lines with colons (likely definitions)
            if 3 < len(line) < 100 and ":" not in line:
                topics.append(line)
        
        # Extract key phrases using simple heuristics
        potential_phrases = re.findall(r'[A-Z][a-z]+(?:\s+[a-z]+){1,3}', content)
        for phrase in potential_phrases:
            if 10 < len(phrase) < 50:  # Reasonable length for a topic
                topics.append(phrase)
    
    # Remove duplicates and very short topics
    return list(set([t for t in topics if len(t) > 3]))

def generate_qa_pairs_simple(content_items, topics, num_samples):
    """Generate simple Q&A pairs using templates."""
    qa_pairs = []
    
    # Make sure we have enough topics
    if len(topics) < 5:
        topics.extend(["this concept", "this topic", "this method", "this approach", "this technique"])
    
    while len(qa_pairs) < num_samples:
        # Select a random content piece
        content_item = random.choice(content_items)
        content_text = content_item["content"]
        
        # Generate a question
        template = random.choice(QUESTION_TEMPLATES)
        
        # Handle different template types
        if template.count("{}") == 1:
            topic = random.choice(topics)
            question = template.format(topic)
        elif template.count("{}") == 2:
            topic1 = random.choice(topics)
            topic2 = random.choice([t for t in topics if t != topic1])
            question = template.format(topic1, topic2)
        
        # Create a simple answer from content
        if content_item.get("type") == "word":
            answer = f"According to the course materials in paragraph {content_item['para_num']} of {content_item['file']}:\n\n"
        elif content_item.get("type") == "text":
            answer = f"According to the course materials in paragraph {content_item['para_num']} of {content_item['file']}:\n\n"
        elif content_item.get("type") == "pdf":
            answer = f"According to the course materials on page {content_item['page_num']}, paragraph {content_item['para_num']} of {content_item['file']}:\n\n"
        else:  # PowerPoint slide
            answer = f"According to the course materials on slide {content_item['slide_num']} of {content_item['file']}:\n\n"
        
        # Format the answer by taking relevant portions of the content
        # If the content is very long, take only a portion
        if len(content_text) > 500:
            paragraphs = content_text.split("\n\n")
            if len(paragraphs) > 2:
                # Take a few paragraphs
                selected = paragraphs[:3]
                answer += "\n\n".join(selected)
            else:
                # Take the first part
                answer += content_text[:500] + "..."
        else:
            answer += content_text
        
        # Record the source
        if content_item.get("type") == "word":
            source = f"{content_item['file']} (Paragraph {content_item['para_num']})"
        elif content_item.get("type") == "text":
            source = f"{content_item['file']} (Paragraph {content_item['para_num']})"
        elif content_item.get("type") == "pdf":
            source = f"{content_item['file']} (Page {content_item['page_num']}, Paragraph {content_item['para_num']})"
        else:  # PowerPoint slide
            source = f"{content_item['file']} (Slide {content_item['slide_num']})"
        
        qa_pairs.append({
            "question": question,
            "answer": answer,
            "source": source
        })
    
    return qa_pairs

def save_qa_pair_immediately(qa_pair, output_file, format_type="finetune"):
    """Save a single Q&A pair immediately to avoid data loss."""
    # Format the pair
    if format_type == "alpaca":
        formatted_data = {
            "instruction": qa_pair["question"],
            "input": "",  # No additional input
            "output": qa_pair["answer"]
        }
    elif format_type == "sharegpt":
        conversation = [
            {"from": "human", "value": qa_pair["question"]},
            {"from": "gpt", "value": qa_pair["answer"]}
        ]
        formatted_data = {"conversations": conversation}
    else:  # "finetune" format for prompt/completion
        formatted_data = {
            "prompt": qa_pair["question"],
            "completion": qa_pair["answer"]
        }
    
    # Append to the output file
    with jsonlines.open(output_file, mode='a') as writer:
        writer.write(formatted_data)

    return True

def generate_qa_pairs_with_ollama(content_items, num_samples, ollama_model, batch_size=5, truly_random=False, 
                                 incremental_save=False, output_file=None, format_type="finetune", start_from=0):
    """Generate Q&A pairs using Ollama model."""
    qa_pairs = []
    processed_count = 0
    saved_count = 0
    retry_limit = 3
    
    print(f"Generating {num_samples} Q&A pairs using Ollama model '{ollama_model}'...")
    print(f"Using batch size of {batch_size}")
    
    if incremental_save and output_file:
        print(f"Incremental saving enabled. Results will be saved to {output_file}")
        
        # Create/clear the file if we're starting from the beginning
        if start_from == 0:
            with open(output_file, 'w') as f:
                pass  # Create empty file or clear if exists
    
    # Check if Ollama is running
    try:
        response = requests.get("http://localhost:11434/api/version")
        if response.status_code != 200:
            print("Warning: Ollama server doesn't seem to be responding correctly.")
    except requests.exceptions.ConnectionError:
        print("Error: Cannot connect to Ollama server. Please make sure Ollama is running.")
        print("You can start Ollama by running 'ollama serve' in a terminal.")
        return []
    
    # Get available models
    try:
        response = requests.get(f"http://localhost:11434/api/tags")
        models_list = "unknown"
        if response.status_code == 200:
            available_models = [model["name"] for model in response.json().get("models", [])]
            models_list = ", ".join(available_models)
            if ollama_model not in available_models:
                print(f"Warning: Model '{ollama_model}' not found in Ollama. Available models: {models_list}")
                print(f"Attempting to use '{ollama_model}' anyway...")
        else:
            print(f"Could not retrieve models list. Status code: {response.status_code}")
    except Exception as e:
        print(f"Warning: Could not verify available models: {e}")
    
    # Prepare content items
    random.shuffle(content_items)
    selected_items = content_items[:num_samples]
    
    # Set item batch size based on model constraints
    if "deepseek" in ollama_model.lower() or "llama3" in ollama_model.lower():
        batch_size = 1  # These models work better with single items
        print(f"Adjusted batch size to 1 for {ollama_model}")
    
    # Process in batches
    for batch_start in range(0, min(len(selected_items), num_samples), batch_size):
        if len(qa_pairs) >= num_samples:
            break
            
        batch_end = min(batch_start + batch_size, len(selected_items), num_samples)
        batch_items = selected_items[batch_start:batch_end]
        
        print(f"Processing batch {batch_start//batch_size + 1} ({batch_start} to {batch_end-1})...")
        
        # Create a simpler prompt - tell the model exactly what we need
        batch_prompt = "You are an education expert creating a Q&A dataset for students. "
        batch_prompt += "For each course material, create clear questions students might ask and detailed answers.\n\n"
        
        # Add content from each item in the batch
        for i, item in enumerate(batch_items):
            content_text = item["content"]
            
            # Determine content type
            if item.get("type") == "word":
                content_type = f"paragraph {item['para_num']} from document"
                source = f"{item['file']} (Paragraph {item['para_num']})"
            elif item.get("type") == "text":
                content_type = f"paragraph {item['para_num']} from text file"
                source = f"{item['file']} (Paragraph {item['para_num']})"
            elif item.get("type") == "pdf":
                content_type = f"page {item['page_num']}, paragraph {item['para_num']} from PDF"
                source = f"{item['file']} (Page {item['page_num']}, Paragraph {item['para_num']})"
            else:  # PowerPoint slide
                content_type = f"slide {item['slide_num']}"
                source = f"{item['file']} (Slide {item['slide_num']})"
                
            batch_prompt += f"\n=== MATERIAL ===\n"
            batch_prompt += f"SOURCE: {source}\n"
            batch_prompt += f"CONTENT:\n{content_text}\n"
            batch_prompt += f"=== END MATERIAL ===\n"
        
        # Add instructions for formatting - make it super clear
        batch_prompt += """
Create 3 different question-answer pairs based on the material above.
Format exactly as follows:

QUESTION 1: [Write a clear student question]
ANSWER 1: [Write a detailed answer]

QUESTION 2: [Write a clear student question]
ANSWER 2: [Write a detailed answer]

QUESTION 3: [Write a clear student question]
ANSWER 3: [Write a detailed answer]

Important:
1. Focus on key concepts from the material
2. Make questions mimic how students actually ask
3. Answer accurately based ONLY on the given content
4. Maintain QUESTION X: and ANSWER X: labels exactly as shown
5. Make answers thorough and educational
"""
        
        # Try to get a response from Ollama
        batch_qa_pairs = []  # Store pairs from this batch
        
        for retry in range(retry_limit):
            try:
                response = requests.post(
                    "http://localhost:11434/api/generate",
                    json={
                        "model": ollama_model,
                        "prompt": batch_prompt,
                        "stream": False,
                        "temperature": 0.3  # Lower temperature for more consistent formatting
                    },
                    timeout=180  # 3-minute timeout for batch processing
                )
                
                if response.status_code == 200:
                    generated_text = response.json().get("response", "")
                    print(f"Got {len(generated_text)} characters of response")
                    
                    # Try multiple parsing approaches for flexibility
                    success = False
                    
                    # Approach 1: Direct QUESTION/ANSWER pattern matching
                    pattern = r"QUESTION\s*(\d+):\s*(.*?)(?:\n|$)ANSWER\s*\1:\s*(.*?)(?:\n\s*QUESTION\s*\d+:|$)"
                    matches = re.findall(pattern, generated_text, re.DOTALL | re.IGNORECASE)
                    
                    if matches:
                        success = True
                        for _, question, answer in matches:
                            batch_qa_pairs.append({
                                "question": question.strip(),
                                "answer": answer.strip(),
                                "source": source
                            })
                    
                    # Approach 2: Try alternative formats if model didn't follow instructions exactly
                    if not success:
                        alt_patterns = [
                            # Q1/A1 format
                            r"Q(?:UESTION)?\s*(\d+):\s*(.*?)(?:\n|$)A(?:NSWER)?\s*\1:\s*(.*?)(?:\n\s*Q(?:UESTION)?\s*\d+:|$)",
                            # Question 1/Answer 1 format
                            r"Question\s*(\d+):\s*(.*?)(?:\n|$)Answer\s*\1:\s*(.*?)(?:\n\s*Question\s*\d+:|$)",
                            # Numbered list: 1. Question / 1. Answer format
                            r"(?<!\d)(\d+)[\.\)]\s+(.*?)(?:\n|$)(?:(?:\1|Answer)[\.\)]\s+)(.*?)(?:\n\s*(?!\1)(?:\d+|Question)[\.\)]|$)"
                        ]
                        
                        for pattern in alt_patterns:
                            matches = re.findall(pattern, generated_text, re.DOTALL | re.IGNORECASE)
                            if matches:
                                success = True
                                for _, question, answer in matches:
                                    batch_qa_pairs.append({
                                        "question": question.strip(),
                                        "answer": answer.strip(),
                                        "source": source
                                    })
                                break
                    
                    # Approach 3: Extract any paragraphs that look like questions and answers
                    if not success:
                        print("Using paragraph analysis to find Q&A pairs...")
                        paragraphs = [p.strip() for p in generated_text.split('\n\n') if p.strip()]
                        
                        # Look for question-like and answer-like paragraphs
                        current_q = None
                        
                        for para in paragraphs:
                            # Check if it's a question (ends with ? or starts with typical question words)
                            if para.endswith('?') or any(para.lower().startswith(q) for q in ['what', 'why', 'how', 'when', 'where', 'who', 'can', 'could']):
                                # Save previous Q&A pair if any
                                if current_q is not None and len(paragraphs) > 1:
                                    batch_qa_pairs.append({
                                        "question": current_q,
                                        "answer": paragraphs[paragraphs.index(current_q) + 1] if paragraphs.index(current_q) < len(paragraphs) - 1 else "",
                                        "source": source
                                    })
                                    
                                current_q = para
                            
                        if batch_qa_pairs:
                            success = True
                    
                    # Last resort - try to find any text with a question mark as a potential question
                    if not success:
                        print("Using question mark detection to find Q&A pairs...")
                        lines = [line.strip() for line in generated_text.split('\n') if line.strip()]
                        
                        for i, line in enumerate(lines):
                            if '?' in line and i < len(lines) - 1:
                                answer_text = lines[i+1]
                                if len(answer_text) > 20:  # Only use if answer has reasonable length
                                    batch_qa_pairs.append({
                                        "question": line.strip(),
                                        "answer": answer_text.strip(),
                                        "source": source
                                    })
                
                    print(f"Extracted {len(batch_qa_pairs)} Q&A pairs from this batch")
                    
                    # The key change is here - immediately save each pair as it's created
                    if incremental_save and output_file:
                        successful_saves = 0
                        for pair in batch_qa_pairs:
                            save_success = save_qa_pair_immediately(pair, output_file, format_type)
                            if save_success:
                                successful_saves += 1
                                saved_count += 1
                        if successful_saves > 0:
                            print(f"Saved {successful_saves} new pairs immediately, total saved: {saved_count}")
                                
                    # If we successfully got any pairs, break retry loop
                    if batch_qa_pairs:
                        break
                    
                elif response.status_code == 404:
                    print(f"Error: Model '{ollama_model}' not found. Available models: {models_list}")
                    print("Try one of these command examples:")
                    print("  ollama pull llama2")
                    print("  ollama pull mistral")
                    print("  ollama pull gemma:7b")
                    return []
                else:
                    print(f"Warning: Ollama returned status code {response.status_code}")
                    if retry == retry_limit - 1:
                        print("Response content:", response.text[:200])
            
            except Exception as e:
                print(f"Error in batch {batch_start//batch_size + 1}, retry {retry+1}: {e}")
                
                if retry == retry_limit - 1:
                    print(f"Failed to process batch after {retry_limit} retries")
        
        # Add the new pairs to our collection
        qa_pairs.extend(batch_qa_pairs)
        processed_count += len(batch_items)
        
        print(f"Generated {len(qa_pairs)} Q&A pairs so far...")
        
        # Break if we've reached the target
        if len(qa_pairs) >= num_samples:
            break
    
    # If we didn't generate any Q&A pairs, provide helpful error message
    if not qa_pairs:
        print("\nOllama failed to generate any Q&A pairs. Common issues:")
        print("1. Model not compatible with the API format (try 'llama2' or 'mistral')")
        print("2. Limited context window (reduce batch_size to 1)")
        print("3. Confused by complex materials (simplify input)")
        print("\nFalling back to template generation...")
        
        # Get topics and use template generation as fallback
        topics = extract_topics_from_content(content_items)
        return generate_qa_pairs_simple(content_items, topics, num_samples)
    
    # Truncate to the requested number of samples
    qa_pairs = qa_pairs[:num_samples]
    
    print(f"Successfully generated {len(qa_pairs)} Q&A pairs using Ollama model '{ollama_model}'")
    return qa_pairs

def format_for_fine_tuning(qa_pairs, format_type):
    """Format the Q&A pairs for fine-tuning."""
    if format_type == "alpaca":
        formatted_data = []
        for pair in qa_pairs:
            formatted_data.append({
                "instruction": pair["question"],
                "input": "",  # No additional input
                "output": pair["answer"]
            })
        return formatted_data
        
    elif format_type == "sharegpt":
        formatted_data = []
        for pair in qa_pairs:
            conversation = [
                {"from": "human", "value": pair["question"]},
                {"from": "gpt", "value": pair["answer"]}
            ]
            formatted_data.append({"conversations": conversation})
        return formatted_data
    
    else:  # "finetune" format for prompt/completion
        formatted_data = []
        for pair in qa_pairs:
            formatted_data.append({
                "prompt": pair["question"],
                "completion": pair["answer"]
            })
        return formatted_data

def save_jsonl(data, filename):
    """Save data as JSONL file."""
    directory = "../data/processed"
    if not os.path.exists(directory):
        os.makedirs(directory)
        
    filepath = os.path.join(directory, filename)
    with open(filepath, 'w', encoding='utf-8') as f:
        for item in data:
            f.write(json.dumps(item) + '\n')
    return filepath

def main():
    """Main function to run the course material Q&A generator."""
    args = parse_args()
    
    print(f"Processing course material files in {args.slides_dir}...")
    
    # Find all supported files
    ppt_files = []
    docx_files = []
    txt_files = []
    pdf_files = []
    
    for file in os.listdir(args.slides_dir):
        file_lower = file.lower()
        file_path = os.path.join(args.slides_dir, file)
        
        if file_lower.endswith('.pptx'):
            ppt_files.append(file_path)
        elif file_lower.endswith('.docx'):
            docx_files.append(file_path)
        elif file_lower.endswith('.txt'):
            txt_files.append(file_path)
        elif file_lower.endswith('.pdf'):
            pdf_files.append(file_path)
    
    if not (ppt_files or docx_files or txt_files or pdf_files):
        print(f"No supported files found in {args.slides_dir}")
        print("Supported file types: .pptx, .docx, .pdf, .txt")
        return
    
    print(f"Found {len(ppt_files)} PowerPoint files, {len(docx_files)} Word documents, {len(pdf_files)} PDF files, and {len(txt_files)} text files")
    
    # Extract content from all files
    all_content = []
    
    # Process PowerPoint files
    for ppt_file in ppt_files:
        slides = extract_slide_content(ppt_file)
        all_content.extend(slides)
    
    # Process Word documents
    for docx_file in docx_files:
        try:
            paragraphs = extract_word_content(docx_file)
            all_content.extend(paragraphs)
        except Exception as e:
            print(f"Error processing Word document {docx_file}: {e}")
    
    # Process PDF files
    for pdf_file in pdf_files:
        try:
            paragraphs = extract_pdf_content(pdf_file)
            all_content.extend(paragraphs)
        except Exception as e:
            print(f"Error processing PDF file {pdf_file}: {e}")
    
    # Process text files
    for txt_file in txt_files:
        try:
            paragraphs = extract_text_content(txt_file)
            all_content.extend(paragraphs)
        except Exception as e:
            print(f"Error processing text file {txt_file}: {e}")
    
    if not all_content:
        print("No content could be extracted from the provided files")
        return
    
    print(f"Extracted content from {len(all_content)} slides/paragraphs across all documents")
    
    # Generate Q&A pairs using Ollama
    if args.use_ollama:
        qa_pairs = generate_qa_pairs_with_ollama(
            all_content, 
            args.num_samples, 
            args.ollama_model,
            args.batch_size,
            incremental_save=True,
            output_file=args.output,
            format_type=args.format
        )
    else:
        # Fall back to template-based generation if not using Ollama
        topics = extract_topics_from_content(all_content)
        print(f"Identified {len(topics)} potential topics")
        print(f"Generating {args.num_samples} Q&A pairs using templates...")
        qa_pairs = generate_qa_pairs_simple(all_content, topics, args.num_samples)
    
    # Format for fine-tuning
    formatted_data = format_for_fine_tuning(qa_pairs, args.format)
    
    # Save the data
    if args.split:
        # Shuffle the data
        random.shuffle(formatted_data)
        
        # Split data 7:2:1
        train_size = int(0.7 * len(formatted_data))
        valid_size = int(0.2 * len(formatted_data))
        
        train_data = formatted_data[:train_size]
        valid_data = formatted_data[train_size:train_size+valid_size]
        test_data = formatted_data[train_size+valid_size:]
        
        print(f"Splitting {len(formatted_data)} examples into: {len(train_data)} train, {len(valid_data)} validation, {len(test_data)} test")
        
        # Save to separate files
        train_file = save_jsonl(train_data, "Train.jsonl")
        valid_file = save_jsonl(valid_data, "Valid.jsonl")
        test_file = save_jsonl(test_data, "Test.jsonl")
        
        print(f"Data split and saved to:")
        print(f"  - Training: {train_file} ({len(train_data)} examples)")
        print(f"  - Validation: {valid_file} ({len(valid_data)} examples)")
        print(f"  - Test: {test_file} ({len(test_data)} examples)")
    else:
        # Write to single output file
        if args.output.endswith('.jsonl'):
            with jsonlines.open(args.output, mode='w') as writer:
                writer.write_all(formatted_data)
        else:
            with open(args.output, 'w', encoding='utf-8') as f:
                json.dump(formatted_data, f, indent=2, ensure_ascii=False)
        
        print(f"Generated {len(formatted_data)} Q&A pairs written to {args.output}")

if __name__ == "__main__":
    main()